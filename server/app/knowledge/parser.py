# app/knowledge/parser.py — 文档解析管道
"""
支持: PDF(pymupdf), Word(python-docx), Excel(openpyxl), PPT, TXT, 图片OCR
输出: 统一的文本块列表，带元数据
"""
import os
import re
import json
import hashlib
from typing import List, Dict, Any, Optional
from datetime import datetime


class DocumentParser:
    """多格式文档解析器，输出结构化文本块"""

    # 支持的文件类型
    MIME_MAP = {
        ".pdf": "pdf", ".docx": "docx", ".doc": "docx",
        ".xlsx": "xlsx", ".xls": "xlsx",
        ".pptx": "pptx", ".ppt": "pptx",
        ".txt": "txt", ".md": "txt",
        ".csv": "csv",
        ".png": "image", ".jpg": "image", ".jpeg": "image",
        ".bmp": "image", ".gif": "image",
    }

    def __init__(self, upload_dir: str = None):
        self.upload_dir = upload_dir or os.path.join(
            os.path.dirname(os.path.dirname(__file__)), "..", "uploads"
        )

    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        解析单个文件，返回结构化结果

        Returns:
            {"filename": str, "type": str, "pages": int, "chunks": [...],
             "metadata": {...}, "full_text": str, "hash": str}
        """
        ext = os.path.splitext(file_path)[1].lower()
        file_type = self.MIME_MAP.get(ext, "unknown")
        filename = os.path.basename(file_path)
        file_hash = self._hash_file(file_path)
        size_kb = os.path.getsize(file_path) / 1024

        if file_type == "pdf":
            text, pages = self._parse_pdf(file_path)
        elif file_type == "docx":
            text, pages = self._parse_docx(file_path)
        elif file_type == "xlsx":
            text, pages = self._parse_xlsx(file_path)
        elif file_type == "pptx":
            text, pages = self._parse_pptx(file_path)
        elif file_type in ("txt", "csv"):
            text, pages = self._parse_text(file_path)
        elif file_type == "image":
            text, pages = self._parse_image(file_path)
        else:
            raise ValueError(f"不支持的文件类型: {ext}")

        if not text or len(text.strip()) < 10:
            raise ValueError(f"文件内容为空或过短 ({len(text)} 字符)")

        chunks = self._chunk_text(text, filename)
        if not chunks:
            raise ValueError("文本分块结果为空")

        return {
            "filename": filename, "file_path": file_path,
            "type": file_type, "pages": pages, "chunks": chunks,
            "full_text": text, "hash": file_hash,
            "size_kb": round(size_kb, 1),
            "parsed_at": datetime.now().isoformat(),
        }

    # ── 各格式解析器 ──────────────────────────────────────────

    def _parse_pdf(self, path: str) -> tuple:
        """PDF解析，返回 (文本, 页数)"""
        try:
            import fitz  # pymupdf
            doc = fitz.open(path)
            pages = []
            for page in doc:
                pages.append(page.get_text("text"))
            doc.close()
            return "\n\n".join(pages), len(pages)
        except ImportError:
            raise ImportError("pymupdf 未安装: pip install pymupdf")

    def _parse_docx(self, path: str) -> tuple:
        """Word解析"""
        try:
            from docx import Document
            doc = Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # 也提取表格
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        paragraphs.append(" | ".join(cells))
            return "\n\n".join(paragraphs), 1
        except ImportError:
            raise ImportError("python-docx 未安装")

    def _parse_xlsx(self, path: str) -> tuple:
        """Excel解析"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            all_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                all_text.append(f"--- {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) if c is not None else "" for c in row]
                    if any(v.strip() for v in vals):
                        all_text.append(" | ".join(vals))
            return "\n".join(all_text), len(wb.sheetnames)
        except ImportError:
            raise ImportError("openpyxl 未安装")

    def _parse_pptx(self, path: str) -> tuple:
        """PPT解析（简化版）"""
        text = []
        try:
            from pptx import Presentation
            prs = Presentation(path)
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
            return "\n\n".join(text), len(prs.slides)
        except ImportError:
            # 降级: 纯文本提取
            return f"[PPT文件: {os.path.basename(path)}，需安装python-pptx]", 0

    def _parse_text(self, path: str) -> tuple:
        """纯文本解析"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        for enc in encodings:
            try:
                with open(path, 'r', encoding=enc) as f:
                    return f.read(), 1
            except UnicodeDecodeError:
                continue
        return f"[无法解码: {path}]", 0

    def _parse_image(self, path: str) -> tuple:
        """图片OCR（占位）"""
        return f"[图片文件: {os.path.basename(path)}，OCR待接入]", 1

    # ── 文本分块 ──────────────────────────────────────────────

    def _chunk_text(self, text: str, source: str = "",
                    chunk_size: int = 800, overlap: int = 100) -> List[Dict]:
        """
        语义分块：按段落+字符数分块，保留重叠用于上下文连贯
        """
        if not text or len(text) < 50:
            return []

        chunks = []
        paragraphs = re.split(r'\n\s*\n', text)

        current_chunk = ""
        chunk_idx = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            if len(current_chunk) + len(para) > chunk_size and current_chunk:
                chunks.append({
                    "id": f"{source}_chunk{chunk_idx}",
                    "text": current_chunk.strip(),
                    "source": source,
                    "index": chunk_idx,
                })
                chunk_idx += 1
                # 保留重叠
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_text + "\n\n" + para
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para

        # 最后一个块
        if current_chunk.strip():
            chunks.append({
                "id": f"{source}_chunk{chunk_idx}",
                "text": current_chunk.strip(),
                "source": source,
                "index": chunk_idx,
            })

        return chunks

    # ── 工具 ──────────────────────────────────────────────────

    def _hash_file(self, path: str) -> str:
        """文件MD5"""
        h = hashlib.md5()
        with open(path, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b''):
                h.update(chunk)
        return h.hexdigest()

    def get_page_count(self, file_path: str) -> int:
        """快速获取文件页数/大小"""
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == '.pdf':
                import fitz
                doc = fitz.open(file_path)
                pages = len(doc)
                doc.close()
                return pages
            return 1
        except:
            return 0
